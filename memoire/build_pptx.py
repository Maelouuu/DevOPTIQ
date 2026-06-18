# -*- coding: utf-8 -*-
"""Génère le PowerPoint de soutenance — De la maquette au produit (Maël Girardin / A.F.D.E.C)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image
import os

A = "_assets"
def asset(p): return os.path.join(A, p)

# ---------- palette ----------
INK   = RGBColor(0x0F,0x17,0x2A)
SLATE = RGBColor(0x47,0x55,0x69)
MUTE  = RGBColor(0x64,0x74,0x8B)
LINE  = RGBColor(0xE2,0xE8,0xF0)
WHITE = RGBColor(0xFF,0xFF,0xFF)
BG    = RGBColor(0xFF,0xFF,0xFF)
PANEL = RGBColor(0xF1,0xF5,0xF9)
PANEL2= RGBColor(0xF8,0xFA,0xFC)
INDIGO= RGBColor(0x43,0x38,0xCA)
INDIGO2=RGBColor(0x63,0x66,0xF1)
INDIGO_LT=RGBColor(0xEE,0xF2,0xFF)
PINK  = RGBColor(0xEC,0x48,0x99)
PINK_DK=RGBColor(0xBE,0x18,0x5D)
PINK_LT=RGBColor(0xFD,0xF2,0xF8)
GREEN = RGBColor(0x22,0xC5,0x5E)
GREEN_DK=RGBColor(0x15,0x80,0x3D)
GREEN_LT=RGBColor(0xF0,0xFD,0xF4)
RED   = RGBColor(0xB9,0x1C,0x1C)
RED_LT= RGBColor(0xFE,0xF2,0xF2)
AMBER = RGBColor(0xB4,0x53,0x09)

TITLE="Georgia"; BODY="Calibri"; MONO="Consolas"

BASE  = dict(afdec=False, accent=INDIGO, kicker=INDIGO, soft=INDIGO_LT, accent2=SLATE)
AFDEC = dict(afdec=True,  accent=PINK,   kicker=PINK_DK, soft=PINK_LT,  accent2=GREEN_DK)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

# ---------- low-level helpers ----------
def shadow_off(sp):
    sp.shadow.inherit = False

def solid(sp, rgb):
    sp.fill.solid(); sp.fill.fore_color.rgb = rgb

def noline(sp):
    sp.line.fill.background()

def setline(sp, rgb, w=1.0):
    sp.line.color.rgb = rgb; sp.line.width = Pt(w)

def rect(slide, x,y,w,h, fill=None, lc=None, lw=1.0):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x,y,w,h); shadow_off(sp)
    if fill is None: sp.fill.background()
    else: solid(sp, fill)
    if lc is None: noline(sp)
    else: setline(sp, lc, lw)
    return sp

def rounded(slide, x,y,w,h, fill=None, lc=None, lw=1.0, rad=0.08):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x,y,w,h); shadow_off(sp)
    try: sp.adjustments[0] = rad
    except Exception: pass
    if fill is None: sp.fill.background()
    else: solid(sp, fill)
    if lc is None: noline(sp)
    else: setline(sp, lc, lw)
    return sp

def oval(slide, x,y,w,h, fill, lc=None, lw=1.0):
    sp = slide.shapes.add_shape(MSO_SHAPE.OVAL, x,y,w,h); shadow_off(sp)
    solid(sp, fill)
    if lc is None: noline(sp)
    else: setline(sp, lc, lw)
    return sp

def _runs(p, content, font, size, color, bold, italic):
    items = content if isinstance(content, list) else [(content, {})]
    for t,o in items:
        r = p.add_run(); r.text = t
        r.font.name = o.get('font', font); r.font.size = Pt(o.get('size', size))
        r.font.bold = o.get('bold', bold); r.font.italic = o.get('italic', italic)
        r.font.color.rgb = o.get('color', color)
        if 'spc' in o:
            r.font._rPr.set('spc', str(o['spc']))

def text(slide, x,y,w,h, lines, font=BODY, size=14, color=INK, bold=False, italic=False,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, sp_after=4, sp_before=0, ls=1.04):
    tb = slide.shapes.add_textbox(x,y,w,h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(sp_after); p.space_before = Pt(sp_before)
        if ls: p.line_spacing = ls
        _runs(p, ln, font, size, color, bold, italic)
    return tb

def bullets(slide, x,y,w,h, items, marker="—", mcolor=INDIGO, color=INK, size=14,
            bold_lead=False, sp_after=8, ls=1.06, anchor=MSO_ANCHOR.TOP):
    lines = []
    for it in items:
        if isinstance(it, tuple):  # (lead, rest)
            lead, rest = it
            lines.append([(marker+"  ", {'color':mcolor,'bold':True}),
                          (lead, {'color':color,'bold':True}),
                          (rest, {'color':color})])
        else:
            lines.append([(marker+"  ", {'color':mcolor,'bold':True}), (it, {'color':color})])
    return text(slide, x,y,w,h, lines, size=size, sp_after=sp_after, ls=ls, anchor=anchor)

def img_fit(slide, path, x,y,maxw,maxh, frame=True, fc=LINE, center=True):
    iw,ih = Image.open(path).size; ar = iw/ih
    w = maxw; h = int(w/ar)
    if h > maxh: h = maxh; w = int(h*ar)
    px = x + (maxw-w)//2 if center else x
    py = y + (maxh-h)//2 if center else y
    if frame:
        fr = rounded(slide, px-Emu(34000), py-Emu(34000), w+Emu(68000), h+Emu(68000),
                     fill=WHITE, lc=fc, lw=1.0, rad=0.05); shadow_off(fr)
    slide.shapes.add_picture(path, px, py, width=w)
    return (px,py,w,h)

# ---------- skeleton ----------
def new_slide(theme, number=None):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = BG
    # left ribbon
    rib = rect(s, 0,0, Inches(0.18), SH, fill=theme['accent'])
    if theme['afdec']:
        rect(s, 0,0, Inches(0.18), Inches(1.15), fill=GREEN)  # green chip (separate element)
        # brand chip top-right
        bx, bw = Inches(10.86), Inches(2.0)
        chip = rounded(s, bx, Inches(0.34), bw, Inches(0.44), fill=PINK_DK, rad=0.5)
        text(s, bx, Inches(0.34), bw-Inches(0.18), Inches(0.44),
             [[("● ", {'color':GREEN,'size':10}), ("A.F.D.E.C", {'color':WHITE,'size':12.5,'bold':True,'spc':120})]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # footer
    text(s, Inches(0.55), Inches(7.04), Inches(8.5), Inches(0.34),
         [[("De la maquette au produit", {'color':SLATE,'size':9,'bold':True}),
           ("   ·   Soutenance — Maël Girardin · A.F.D.E.C", {'color':MUTE,'size':9})]],
         anchor=MSO_ANCHOR.MIDDLE)
    if number is not None:
        text(s, Inches(12.1), Inches(7.0), Inches(0.95), Inches(0.4),
             [[(f"{number:02d}", {'color':theme['accent'],'size':13,'bold':True})]],
             align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    return s

def title_block(s, kicker, title, theme, ty=Inches(0.6), tsize=30, tw=Inches(11.6)):
    text(s, Inches(0.64), ty, tw, Inches(0.3),
         [[(kicker, {'color':theme['kicker'],'size':12.5,'bold':True,'spc':190})]])
    text(s, Inches(0.6), ty+Inches(0.33), tw, Inches(1.0),
         [[(title, {'font':TITLE,'color':INK,'size':tsize,'bold':True})]], ls=1.0)
    rect(s, Inches(0.66), ty+Inches(1.16), Inches(1.5), Pt(3.2), fill=theme['accent'])

# ---------- composite components ----------
def stat(s, x,y,w,h, big, label, accent, fill=PANEL2):
    rounded(s, x,y,w,h, fill=fill, lc=LINE, lw=1.0, rad=0.10)
    rect(s, x+Inches(0.0), y, Inches(0.07), h, fill=accent)
    text(s, x+Inches(0.2), y+Inches(0.12), w-Inches(0.3), Inches(0.5),
         [[(big, {'font':TITLE,'color':accent,'size':25,'bold':True})]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, x+Inches(0.2), y+h-Inches(0.62), w-Inches(0.3), Inches(0.55),
         [[(label, {'color':SLATE,'size':10.5})]], anchor=MSO_ANCHOR.MIDDLE, ls=1.0)

def pill(s, x,y,w,h, label, fill, tc, size=11.5):
    rounded(s, x,y,w,h, fill=fill, rad=0.5)
    text(s, x,y,w,h, [[(label,{'color':tc,'size':size,'bold':True})]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

def step_banner(s, x,y,w, label_lead, label_rest, theme):
    """'Le pas franchi' style banner."""
    h = Inches(0.62)
    rounded(s, x,y,w,h, fill=theme['soft'], lc=None, rad=0.16)
    rect(s, x, y, Inches(0.08), h, fill=theme['accent'])
    text(s, x+Inches(0.24), y, w-Inches(0.45), h,
         [[(label_lead, {'color':theme['kicker'],'size':12,'bold':True}),
           (label_rest, {'color':INK,'size':12.5})]], anchor=MSO_ANCHOR.MIDDLE, ls=1.02)

def code_card(s, x,y,w,h, header, hcolor, code, accent_border):
    rounded(s, x,y,w,h, fill=RGBColor(0xFB,0xFD,0xFF), lc=accent_border, lw=1.2, rad=0.05)
    text(s, x+Inches(0.18), y+Inches(0.1), w-Inches(0.3), Inches(0.3),
         [[(header, {'color':hcolor,'size':11,'bold':True})]])
    rect(s, x+Inches(0.18), y+Inches(0.42), w-Inches(0.36), Pt(1), fill=LINE)
    lines = [[(ln, {'font':MONO,'size':9.2,'color':INK})] for ln in code]
    text(s, x+Inches(0.2), y+Inches(0.5), w-Inches(0.34), h-Inches(0.6),
         lines, sp_after=1, ls=1.06)

def takeaway(s, x,y,w, lead, rest, theme):
    h=Inches(0.6)
    rounded(s, x,y,w,h, fill=PANEL2, lc=LINE, lw=1.0, rad=0.12)
    rect(s, x, y, Inches(0.07), h, fill=theme['accent2'])
    text(s, x+Inches(0.22), y, w-Inches(0.4), h,
         [[(lead, {'color':theme['accent2'] if theme['afdec'] else theme['accent'],'size':11.5,'bold':True}),
           (rest, {'color':INK,'size':11.5})]], anchor=MSO_ANCHOR.MIDDLE, ls=1.03)

def column_card(s, x,y,w,h, head, head_color, items, marker, fill=PANEL2):
    rounded(s, x,y,w,h, fill=fill, lc=LINE, lw=1.0, rad=0.06)
    rounded(s, x, y, w, Inches(0.5), fill=head_color, rad=0.06)
    rect(s, x, y+Inches(0.25), w, Inches(0.25), fill=head_color)  # square off bottom of header
    text(s, x+Inches(0.2), y, w-Inches(0.3), Inches(0.5),
         [[(head,{'color':WHITE,'size':13,'bold':True})]], anchor=MSO_ANCHOR.MIDDLE)
    bullets(s, x+Inches(0.22), y+Inches(0.66), w-Inches(0.42), h-Inches(0.8),
            items, marker=marker, mcolor=head_color, size=11.3, sp_after=7, ls=1.05)

# =========================================================================
# SLIDE 1 — COVER (no number)
# =========================================================================
s = prs.slides.add_slide(BLANK)
s.background.fill.solid(); s.background.fill.fore_color.rgb = BG
# top brand bands (separate pink & green, no blend)
rect(s, 0,0, SW, Inches(0.16), fill=PINK)
rect(s, 0, Inches(0.16), SW, Inches(0.05), fill=GREEN)
# logos
img_fit(s, asset("img01.png"), Inches(0.7), Inches(0.6), Inches(2.6), Inches(1.0), frame=False, center=False)
img_fit(s, asset("img00.png"), Inches(10.4), Inches(0.62), Inches(2.2), Inches(0.95), frame=False, center=False)
# kicker
text(s, Inches(0.8), Inches(2.25), Inches(11.7), Inches(0.4),
     [[("MÉMOIRE D'ALTERNANCE  ·  B.U.T. INFORMATIQUE — 3ᵉ ANNÉE  ·  2025–2026",
        {'color':PINK_DK,'size':13,'bold':True,'spc':150})]], align=PP_ALIGN.CENTER)
# big title
text(s, Inches(0.6), Inches(2.75), Inches(12.1), Inches(1.5),
     [[("De la maquette au produit", {'font':TITLE,'color':INK,'size':58,'bold':True})]],
     align=PP_ALIGN.CENTER)
# accent motif (pink rule + green square, separate)
rect(s, Inches(5.66), Inches(4.18), Inches(1.7), Pt(3.4), fill=PINK)
rect(s, Inches(7.42), Inches(4.12), Inches(0.16), Inches(0.16), fill=GREEN)
# subtitle
text(s, Inches(2.0), Inches(4.45), Inches(9.33), Inches(0.9),
     [[("Anatomie d'une professionnalisation : faire grandir une application de management ",
        {'color':SLATE,'size':15,'italic':True})],
      [("des compétences — et grandir avec elle.", {'color':SLATE,'size':15,'italic':True})]],
     align=PP_ALIGN.CENTER, ls=1.15)
# meta band
rounded(s, Inches(1.7), Inches(5.75), Inches(9.93), Inches(1.05), fill=PANEL2, lc=LINE, lw=1.0, rad=0.10)
text(s, Inches(2.0), Inches(5.92), Inches(9.4), Inches(0.4),
     [[("Maël GIRARDIN", {'color':INK,'size':16,'bold':True}),
       ("      —      IUT — site d'Arles  ·  Département Informatique", {'color':SLATE,'size':12.5})]],
     align=PP_ALIGN.CENTER)
text(s, Inches(2.0), Inches(6.34), Inches(9.4), Inches(0.4),
     [[("Maître d'alternance : ", {'color':MUTE,'size':12}), ("Hubert GRANDJEAN", {'color':PINK_DK,'size':12,'bold':True}),
       ("        Tuteur IUT : ", {'color':MUTE,'size':12}), ("Bruno COLOMBEL", {'color':GREEN_DK,'size':12,'bold':True})]],
     align=PP_ALIGN.CENTER)

# =========================================================================
# SLIDE 2 — PLAN
# =========================================================================
s = new_slide(BASE, 2)
title_block(s, "DÉROULÉ DE LA SOUTENANCE", "Plan", BASE)
plan = [
    ("L'entreprise & le projet", "l'A.F.D.E.C, la méthode OPTIQ, et le basculement maquette → produit"),
    ("Le développement", "quatre franchissements dans le code, en regard « avant / après »"),
    ("Professionnaliser le produit", "base de données en production, tests, documentation"),
    ("Analyse critique", "l'arbitrage assumé, la sécurité, le plan de durcissement"),
    ("Formation & entreprise", "ce que l'IUT a apporté, ce que j'ai dû renforcer"),
]
y = Inches(2.15); ph = Inches(0.78); gap = Inches(0.08)
for i,(lead,rest) in enumerate(plan):
    yy = y + i*(ph+gap)
    rounded(s, Inches(0.7), yy, Inches(11.9), ph, fill=PANEL2, lc=LINE, lw=1.0, rad=0.12)
    oval(s, Inches(0.92), yy+Inches(0.16), Inches(0.46), Inches(0.46), fill=INDIGO_LT)
    text(s, Inches(0.92), yy+Inches(0.16), Inches(0.46), Inches(0.46),
         [[(str(i+1), {'color':INDIGO,'size':16,'bold':True})]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(1.6), yy, Inches(10.7), ph,
         [[(lead, {'color':INK,'size':15.5,'bold':True}), ("    —  "+rest, {'color':SLATE,'size':12.5})]],
         anchor=MSO_ANCHOR.MIDDLE, ls=1.0)

# =========================================================================
# SLIDE 3 — PROBLÉMATIQUE
# =========================================================================
s = new_slide(BASE, 3)
title_block(s, "LA QUESTION QUI GUIDE CE TRAVAIL", "Problématique", BASE)
rounded(s, Inches(1.1), Inches(2.5), Inches(11.1), Inches(2.7), fill=INDIGO_LT, lc=None, rad=0.06)
rect(s, Inches(1.1), Inches(2.5), Inches(0.1), Inches(2.7), fill=INDIGO)
text(s, Inches(1.7), Inches(2.5), Inches(10.1), Inches(2.7),
     [[("Comment transformer une ", {'color':INK,'size':23}),
       ("maquette", {'color':PINK_DK,'size':23,'bold':True}),
       (" en un ", {'color':INK,'size':23}),
       ("produit professionnel et commercialisable", {'color':GREEN_DK,'size':23,'bold':True}),
       (", et en quoi cette transformation constitue-t-elle, pour le développeur en alternance, un véritable ",
        {'color':INK,'size':23}),
       ("levier de montée en compétences", {'color':INDIGO,'size':23,'bold':True}),
       (" ?", {'color':INK,'size':23})]],
     anchor=MSO_ANCHOR.MIDDLE, ls=1.22)
text(s, Inches(1.1), Inches(5.5), Inches(11.1), Inches(0.6),
     [[("Deux trajectoires, une même histoire : faire grandir le projet — et grandir avec lui.",
        {'color':SLATE,'size':14,'italic':True})]], align=PP_ALIGN.CENTER)

# =========================================================================
# SLIDE 4 — A.F.D.E.C & MÉTHODE OPTIQ
# =========================================================================
s = new_slide(AFDEC, 4)
title_block(s, "PARTIE I · CONTEXTE", "L'A.F.D.E.C & la méthode OPTIQ", AFDEC)
# left : AFDEC
column_card(s, Inches(0.7), Inches(2.15), Inches(5.7), Inches(2.55), "L'A.F.D.E.C", PINK_DK,
            [("Laboratoire", " de recherche appliquée sur les organisations apprenantes et performantes."),
             ("Management des compétences", ", cartographie du travail réel, pilotage des activités et des rôles."),
             ("Normes françaises", " : X50-124 (compétences) et X50-766 (habiletés socio-cognitives).")],
            marker="—", fill=PINK_LT)
# right : OPTIQ
column_card(s, Inches(6.6), Inches(2.15), Inches(6.0), Inches(2.55), "La méthode OPTIQ©", GREEN_DK,
            [("CartoColor©", " — la carte des activités : chacun se situe dans le flux et voit ses compétences clés."),
             ("OPTIQFluent©", " — compétences, performances, charge de travail ; des fiches de poste aux fiches de rôle."),
             ("La cartographie", " = le point d'entrée central de l'application.")],
            marker="—", fill=GREEN_LT)
# bottom objective
rounded(s, Inches(0.7), Inches(4.95), Inches(11.9), Inches(0.95), fill=PANEL2, lc=LINE, lw=1.0, rad=0.08)
rect(s, Inches(0.7), Inches(4.95), Inches(0.09), Inches(0.95), fill=PINK)
text(s, Inches(1.0), Inches(4.95), Inches(11.3), Inches(0.95),
     [[("L'objectif  ", {'color':PINK_DK,'size':13,'bold':True}),
       ("Rendre la méthode OPTIQ accessible à tous — sans en trahir la finesse : ni base trop rigide, ni interface trop verbeuse.",
        {'color':INK,'size':13.5})]], anchor=MSO_ANCHOR.MIDDLE, ls=1.1)

# =========================================================================
# SLIDE 5 — MON RÔLE
# =========================================================================
s = new_slide(AFDEC, 5)
title_block(s, "PARTIE I · MON RÔLE", "Du stagiaire au co-concepteur", AFDEC)
# progression timeline
steps = [("Stage — 2ᵉ année","12 semaines","Reprendre l'existant, consolider, livrer une maquette à présenter."),
         ("Été — transition","la continuité","L'idée d'un produit complet commence à mûrir."),
         ("Alternance — 3ᵉ année","une année","Transformer la maquette en produit : structure, visuel, tests, doc.")]
x0=Inches(0.7); cw=Inches(3.83); gap=Inches(0.20); y=Inches(2.2)
cols=[PINK, GREEN_DK, PINK_DK]
for i,(t,tag,d) in enumerate(steps):
    x = x0 + i*(cw+gap)
    rounded(s, x, y, cw, Inches(1.85), fill=PANEL2, lc=LINE, lw=1.0, rad=0.08)
    rect(s, x, y, cw, Inches(0.09), fill=cols[i])
    text(s, x+Inches(0.22), y+Inches(0.2), cw-Inches(0.4), Inches(0.4),
         [[(t, {'color':INK,'size':14.5,'bold':True})]])
    pill(s, x+Inches(0.22), y+Inches(0.62), Inches(1.5), Inches(0.33), tag, cols[i] if i!=1 else GREEN_DK, WHITE, size=10.5)
    text(s, x+Inches(0.22), y+Inches(1.05), cw-Inches(0.42), Inches(0.7),
         [[(d, {'color':SLATE,'size':11.5})]], ls=1.08)
    if i<2:
        text(s, x+cw-Inches(0.02), y+Inches(0.7), gap+Inches(0.1), Inches(0.4),
             [[("→",{'color':MUTE,'size':18,'bold':True})]], align=PP_ALIGN.CENTER)
# shift + quote
bullets(s, Inches(0.72), Inches(4.35), Inches(7.4), Inches(2.0),
        [("D'exécutant", " de missions à interlocuteur avec qui l'on réfléchit le produit."),
         ("Co-construction", " : je propose, Hubert challenge, on arbitre — d'égal à égal."),
         ("Une responsabilité de pilotage", " : priorités, séquencement, versions stables pendant les absences.")],
        marker="—", mcolor=PINK, size=13.5, sp_after=11)
rounded(s, Inches(8.35), Inches(4.35), Inches(4.25), Inches(1.95), fill=PINK_LT, lc=None, rad=0.08)
rect(s, Inches(8.35), Inches(4.35), Inches(0.1), Inches(1.95), fill=PINK)
text(s, Inches(8.62), Inches(4.35), Inches(3.85), Inches(1.95),
     [[("« Je ne délivrais plus un travail :", {'color':PINK_DK,'size':16,'bold':True,'font':TITLE})],
      [("je co-soutenais une vision. »", {'color':PINK_DK,'size':16,'bold':True,'font':TITLE})]],
     anchor=MSO_ANCHOR.MIDDLE, ls=1.2)

# =========================================================================
# SLIDE 6 — LE BASCULEMENT
# =========================================================================
s = new_slide(AFDEC, 6)
title_block(s, "PARTIE I · LE TOURNANT", "La maquette devient trop convaincante pour rester une maquette", AFDEC, tsize=25)
chantiers = [("①  Solidifier la structure","Base de données robuste, migrations sûres, gestion d'erreurs : un produit ne plante pas chez le client."),
             ("②  Soigner le visuel","Une identité (rose / vert / blanc) pour passer d'un outil fonctionnel à un produit désirable."),
             ("③  Garantir la fiabilité","Documentation et batterie de tests : un produit rend compte de son propre état de santé."),
             ("④  Penser le multi-client","Comptes et entités pour cloisonner les données : le marqueur d'une logique commerciale.")]
x0=Inches(0.7); cw=Inches(5.82); ch=Inches(1.28); gx=Inches(0.26); gy=Inches(0.2); y0=Inches(2.05)
cc=[PINK_DK, GREEN_DK, GREEN_DK, PINK_DK]
for i,(t,d) in enumerate(chantiers):
    x = x0 + (i%2)*(cw+gx); yy = y0 + (i//2)*(ch+gy)
    rounded(s, x, yy, cw, ch, fill=PANEL2, lc=LINE, lw=1.0, rad=0.08)
    rect(s, x, yy, Inches(0.09), ch, fill=cc[i])
    text(s, x+Inches(0.26), yy+Inches(0.14), cw-Inches(0.45), Inches(0.4),
         [[(t, {'color':INK,'size':15,'bold':True})]])
    text(s, x+Inches(0.26), yy+Inches(0.55), cw-Inches(0.5), Inches(0.7),
         [[(d, {'color':SLATE,'size':11.8})]], ls=1.08)
rounded(s, Inches(0.7), Inches(5.6), Inches(11.9), Inches(0.82), fill=GREEN_LT, lc=None, rad=0.08)
rect(s, Inches(0.7), Inches(5.6), Inches(0.09), Inches(0.82), fill=GREEN_DK)
text(s, Inches(1.0), Inches(5.6), Inches(11.4), Inches(0.82),
     [[("Le pourquoi  ", {'color':GREEN_DK,'size':13,'bold':True}),
       ("On ne présente plus une idée à financer, mais une solution à acheter : l'exercice de dev devient un projet d'ingénierie produit.",
        {'color':INK,'size':13})]], anchor=MSO_ANCHOR.MIDDLE, ls=1.05)

# =========================================================================
# SLIDE 7 — AMPLEUR
# =========================================================================
s = new_slide(AFDEC, 7)
title_block(s, "PARTIE II · CE QUI A ÉTÉ CONSTRUIT", "Une application web complète, ~30 modules", AFDEC)
stats = [("~55","Blueprints Flask",PINK_DK), ("29","Modèles de données",GREEN_DK),
         ("63","Gabarits (templates)",PINK_DK), ("~69k","Lignes de code",GREEN_DK)]
x0=Inches(0.7); cw=Inches(1.92); gap=Inches(0.13); y=Inches(2.15)
for i,(b,l,c) in enumerate(stats):
    stat(s, x0+i*(cw+gap), y, cw, Inches(1.15), b, l, c)
# domains
doms = ["Concepteur de cartographie","Activités & tâches","Compétences (savoirs · HSC)",
        "Gestion RH & comptes","Performance & temps","Briques d'IA"]
text(s, Inches(0.72), Inches(3.55), Inches(8.0), Inches(0.3),
     [[("Les grands domaines fonctionnels", {'color':INK,'size':14,'bold':True})]])
yy=Inches(3.95)
for i,d in enumerate(doms):
    x = Inches(0.7) + (i%2)*Inches(4.05); yr = yy + (i//2)*Inches(0.62)
    rounded(s, x, yr, Inches(3.85), Inches(0.5), fill=PANEL2, lc=LINE, lw=1.0, rad=0.2)
    oval(s, x+Inches(0.2), yr+Inches(0.19), Inches(0.13), Inches(0.13), fill=GREEN if i%2 else PINK)
    text(s, x+Inches(0.5), yr, Inches(3.3), Inches(0.5), [[(d,{'color':INK,'size':11.8})]], anchor=MSO_ANCHOR.MIDDLE)
# visual carto
img_fit(s, asset("img15.jpg"), Inches(8.95), Inches(3.6), Inches(3.65), Inches(2.45), frame=True)
text(s, Inches(8.95), Inches(6.05), Inches(3.65), Inches(0.3),
     [[("Le Concepteur de cartographie, aujourd'hui", {'color':MUTE,'size':9.5,'italic':True})]], align=PP_ALIGN.CENTER)
takeaway(s, Inches(0.7), Inches(6.32), Inches(8.05), "Priorité assumée  ",
         "d'abord la richesse fonctionnelle ; la sécurité de pointe et l'optimisation, plus tard (Partie III).", AFDEC)

# =========================================================================
# SLIDE 8 — CAS 1
# =========================================================================
s = new_slide(AFDEC, 8)
title_block(s, "PARTIE II · CAS N°1 — AUTONOMIE", "Notre propre outil : le Concepteur de cartographie", AFDEC, tsize=26)
step_banner(s, Inches(0.7), Inches(1.95), Inches(11.9),
            "Le pas franchi   ", "Cartographie dessinée sous Visio (outil tiers payant)  →  éditeur SVG intégré, autonome", AFDEC)
img_fit(s, asset("fig1.png"), Inches(0.7), Inches(2.8), Inches(7.4), Inches(3.0), frame=False)
# right: real editor + key facts
img_fit(s, asset("img04.jpg"), Inches(8.4), Inches(2.8), Inches(4.2), Inches(1.95), frame=True)
text(s, Inches(8.4), Inches(4.82), Inches(4.2), Inches(0.3),
     [[("L'éditeur intégré, en vrai", {'color':MUTE,'size':9.5,'italic':True})]], align=PP_ALIGN.CENTER)
bullets(s, Inches(8.42), Inches(5.18), Inches(4.2), Inches(1.0),
        [("> 6 500 lignes", " de JavaScript, entièrement maison."),
         ("Dessin → JSON → base", " : la carte est une structure vivante, pas une image.")],
        marker="—", mcolor=GREEN_DK, size=11.3, sp_after=7)
takeaway(s, Inches(0.7), Inches(6.0), Inches(7.4), "Ce que ce choix traduit  ",
         "supprimer une dépendance tierce, c'est rendre le produit vendable et autosuffisant.", AFDEC)

# =========================================================================
# SLIDE 9 — CAS 2
# =========================================================================
s = new_slide(AFDEC, 9)
title_block(s, "PARTIE II · CAS N°2 — LE VIRAGE COMMERCIAL", "Le système multi-entités", AFDEC, tsize=26)
step_banner(s, Inches(0.7), Inches(1.95), Inches(11.9),
            "Le pas franchi   ", "Les données d'une seule organisation  →  un espace cloisonné et étanche par client", AFDEC)
img_fit(s, asset("fig2.png"), Inches(0.7), Inches(2.8), Inches(7.3), Inches(3.0), frame=False)
# right code + stat
code_card(s, Inches(8.25), Inches(2.82), Inches(4.35), Inches(1.95), "● APRÈS — isolation par propriétaire", GREEN_DK,
          ["active = Entity.get_active()","activities = (Activities.query","   .filter_by(entity_id=active.id)","   .all())   # ← chaque client",""], GREEN)
stat(s, Inches(8.25), Inches(4.95), Inches(2.1), Inches(1.0), "+340", "filtres entity_id", PINK_DK)
stat(s, Inches(10.5), Inches(4.95), Inches(2.1), Inches(1.0), "230", "réf. entité active", GREEN_DK)
takeaway(s, Inches(0.7), Inches(6.05), Inches(7.3), "Selon moi  ",
         "le marqueur le plus net du passage au produit : repenser chaque requête de l'application.", AFDEC)

# =========================================================================
# SLIDE 10 — CAS 3
# =========================================================================
s = new_slide(AFDEC, 10)
title_block(s, "PARTIE II · CAS N°3 — DURABILITÉ", "Penser le cloud : le stockage des fichiers", AFDEC, tsize=26)
step_banner(s, Inches(0.7), Inches(1.95), Inches(11.9),
            "Le pas franchi   ", "Des fichiers écrits sur le disque (volatils en cloud)  →  un stockage durable, en base", AFDEC)
code_card(s, Inches(0.7), Inches(2.85), Inches(5.85), Inches(2.35), "● AVANT — disque du conteneur (volatil)", RED,
          ["UPLOAD_DIR = os.path.join(root, \"uploads\")","","def upload_file():","    f = request.files[\"file\"]","    f.save(path)   # perdu au redémarrage"], RED)
code_card(s, Inches(6.75), Inches(2.85), Inches(5.85), Inches(2.35), "● APRÈS — binaire persistant en base", GREEN_DK,
          ["class FileBlob(db.Model):","    data = db.Column(db.LargeBinary)  # BYTEA","","blob = FileBlob(data=f.read())","db.session.add(blob); db.session.commit()"], GREEN)
rounded(s, Inches(0.7), Inches(5.42), Inches(11.9), Inches(0.78), fill=PINK_LT, lc=None, rad=0.08)
rect(s, Inches(0.7), Inches(5.42), Inches(0.09), Inches(0.78), fill=PINK)
text(s, Inches(1.0), Inches(5.42), Inches(11.4), Inches(0.78),
     [[("Le contexte  ", {'color':PINK_DK,'size':12.5,'bold':True}),
       ("Google Cloud Run : système de fichiers éphémère. Au redémarrage, le disque est effacé — les fichiers des clients s'évanouissent.",
        {'color':INK,'size':12.5})]], anchor=MSO_ANCHOR.MIDDLE, ls=1.05)
takeaway(s, Inches(0.7), Inches(6.32), Inches(11.9), "Ce que ce choix traduit  ",
         "cesser de coder « en général » pour coder pour son environnement réel : comprendre la cible de déploiement avant le code.", AFDEC)

# =========================================================================
# SLIDE 11 — CAS 4
# =========================================================================
s = new_slide(AFDEC, 11)
title_block(s, "PARTIE II · CAS N°4 — LE SOIN DU PRODUIT", "Désirable & increvable", AFDEC, tsize=26)
# left desirable
text(s, Inches(0.72), Inches(1.95), Inches(6.0), Inches(0.35),
     [[("Désirable  ", {'color':PINK_DK,'size':14,'bold':True}), ("— une identité qui donne envie", {'color':SLATE,'size':12})]])
img_fit(s, asset("img05.jpg"), Inches(0.7), Inches(2.4), Inches(2.95), Inches(1.65), frame=True)
img_fit(s, asset("img06.jpg"), Inches(3.75), Inches(2.4), Inches(2.95), Inches(1.65), frame=True)
text(s, Inches(0.7), Inches(4.08), Inches(2.95), Inches(0.3), [[("avant — générique",{'color':MUTE,'size':9.5,'italic':True})]], align=PP_ALIGN.CENTER)
text(s, Inches(3.75), Inches(4.08), Inches(2.95), Inches(0.3), [[("après — identité de marque",{'color':MUTE,'size':9.5,'italic':True})]], align=PP_ALIGN.CENTER)
bullets(s, Inches(0.72), Inches(4.5), Inches(6.0), Inches(0.9),
        [("Un design system", " : jetons de couleur rose #ec4899 / vert #22c55e / blanc, sur toute l'application.")],
        marker="—", mcolor=PINK, size=12, sp_after=4)
# right increvable
text(s, Inches(7.0), Inches(1.95), Inches(5.6), Inches(0.35),
     [[("Increvable  ", {'color':GREEN_DK,'size':14,'bold':True}), ("— une IA qui ne fait jamais planter", {'color':SLATE,'size':12})]])
code_card(s, Inches(6.95), Inches(2.4), Inches(5.65), Inches(1.95), "● APRÈS — client optionnel + repli déterministe", GREEN_DK,
          ["def openai_client_or_none():","    if not os.environ.get(\"OPENAI_API_KEY\"):","        return None, \"clé absente\"  # jamais 500","    ...","# repli si l'IA est indisponible → l'app reste debout"], GREEN)
rounded(s, Inches(6.95), Inches(4.5), Inches(5.65), Inches(0.9),
        fill=RGBColor(0xFF,0xFB,0xEB), lc=RGBColor(0xF5,0xD0,0x90), lw=1.0, rad=0.1)
text(s, Inches(7.2), Inches(4.5), Inches(5.2), Inches(0.9),
     [[("Cadre & éthique  ", {'color':AMBER,'size':11.5,'bold':True}),
       ("l'IA propose, l'humain valide. Confidentialité, traçabilité des prompts, prévention du prompt injection : une doctrine A.F.D.E.C.",
        {'color':INK,'size':11})]], anchor=MSO_ANCHOR.MIDDLE, ls=1.04)
takeaway(s, Inches(0.7), Inches(5.62), Inches(11.9), "Ce que ces finitions traduisent  ",
         "deux exigences qu'une démo ignore : séduire, et ne jamais décevoir — car dans un produit, l'erreur n'est pas permise.", AFDEC)

# =========================================================================
# SLIDE 12 — PROFESSIONNALISER
# =========================================================================
s = new_slide(AFDEC, 12)
title_block(s, "PARTIE II · CHANGEMENT DE RÉGIME", "Professionnaliser : structure, tests, documentation", AFDEC, tsize=25)
col = [("Fiabiliser la BDD", GREEN_DK, GREEN_LT,
        [("PostgreSQL", " en prod, SQLite en dev & tests."),
         ("ALTER TABLE idempotents", " au démarrage — nés d'un incident Alembic (verrou, crash en boucle)."),
         ("Robustesse", " : UPSERT, journal d'activité automatique.")]),
       ("Tester", PINK_DK, PINK_LT,
        [("+ 1 000 tests", " automatisés, 41 fichiers, ~16k lignes."),
         ("/testpanel", " : l'état de santé, lisible depuis l'interface."),
         ("~93 % de réussite", " au dernier run — la peur des régressions devient confiance.")]),
       ("Documenter", GREEN_DK, GREEN_LT,
        [("Doc technique", " : architecture, modèles, API, déploiement."),
         ("Guide utilisateur", " pour les futurs clients."),
         ("Transmettre", " plutôt que détenir : un actif repris, audité, étendu.")])]
x0=Inches(0.7); cw=Inches(3.83); gap=Inches(0.20); y=Inches(2.1)
for i,(h,hc,fl,items) in enumerate(col):
    column_card(s, x0+i*(cw+gap), y, cw, Inches(2.7), h, hc, items, marker="—", fill=fl)
img_fit(s, asset("img07.jpg"), Inches(0.7), Inches(5.05), Inches(5.2), Inches(1.55), frame=True, center=False)
text(s, Inches(6.1), Inches(5.2), Inches(6.5), Inches(1.4),
     [[("Le panneau de tests intégré", {'color':INK,'size':13,'bold':True})],
      [("Taux de réussite, durée moyenne, historique des runs, détail jusqu'au test unitaire. ",
        {'color':SLATE,'size':12})],
      [("L'application rend compte, run après run, de son propre état de santé — y compris pour un non-développeur.",
        {'color':SLATE,'size':12})]], ls=1.12, sp_after=4)

# =========================================================================
# SLIDE 13 — ANALYSE CRITIQUE
# =========================================================================
s = new_slide(BASE, 13)
title_block(s, "PARTIE III · ANALYSE CRITIQUE", "Le fonctionnel d'abord, la sécurité ensuite — un choix assumé", BASE, tsize=25)
column_card(s, Inches(0.7), Inches(2.1), Inches(5.85), Inches(2.5), "✓  Atteint", GREEN_DK,
            [("Produit déployé", ", multi-entités, cloisonnement strict des données."),
             ("Identité visuelle", " cohérente + mots de passe hachés (werkzeug)."),
             ("+1 000 tests", " et documentation technique & utilisateur.")],
            marker="—", fill=GREEN_LT)
column_card(s, Inches(6.75), Inches(2.1), Inches(5.85), Inches(2.5), "→  Encore ouvert", AMBER,
            [("Durcissement sécurité", " : CSRF, rate-limit, garde d'accès systématique, en-têtes HTTP."),
             ("Optimisation fine", " : requêtes, cache (volumes encore modestes)."),
             ("Distribution", " : installer, facturer, accompagner — le défi suivant.")],
            marker="—", fill=RGBColor(0xFF,0xFB,0xEB))
# durcissement plan
text(s, Inches(0.72), Inches(4.85), Inches(8.0), Inches(0.3),
     [[("Un plan de durcissement avant commercialisation", {'color':INK,'size':14,'bold':True})]])
horizons=[("Court terme","tests en conditions réelles, sécurité, transfert des données"),
          ("Moyen terme","sauvegarde/récupération, montée en charge, assurance qualité"),
          ("Long terme","SaaS mature : audit, monitoring, RGPD chez le client")]
x0=Inches(0.7); cw=Inches(3.83); gap=Inches(0.20); y=Inches(5.25)
for i,(t,d) in enumerate(horizons):
    x=x0+i*(cw+gap)
    rounded(s, x, y, cw, Inches(1.05), fill=INDIGO_LT, lc=None, rad=0.1)
    rect(s, x, y, cw, Inches(0.07), fill=INDIGO)
    text(s, x+Inches(0.22), y+Inches(0.14), cw-Inches(0.4), Inches(0.35), [[(t,{'color':INDIGO,'size':13,'bold':True})]])
    text(s, x+Inches(0.22), y+Inches(0.5), cw-Inches(0.42), Inches(0.5), [[(d,{'color':SLATE,'size':11})]], ls=1.06)
text(s, Inches(0.72), Inches(6.45), Inches(11.8), Inches(0.4),
     [[("Ces manques ne sont pas des oublis, mais une ", {'color':SLATE,'size':12,'italic':True}),
       ("dette technique consciente, identifiée et priorisée.", {'color':INDIGO,'size':12,'bold':True,'italic':True})]])

# =========================================================================
# SLIDE 14 — LIAISON IUT <-> ENTREPRISE
# =========================================================================
s = new_slide(BASE, 14)
title_block(s, "FORMATION & ENTREPRISE · REGARDS CROISÉS", "Ce que l'IUT a apporté, ce que j'ai dû renforcer", BASE, tsize=25)
cols14=[("✓  Ce qui a aidé", GREEN_DK, GREEN_LT,
         [("Pile technique", " : HTML, CSS, JS, SQL — familiers depuis la 2ᵉ année web."),
          ("Python", " : proche des langages objet déjà maîtrisés."),
          ("Agilité & comptes rendus", " : directement issus des SAE et projets de groupe.")]),
        ("≈  Ce que l'école n'enseigne qu'en théorie", INDIGO, INDIGO_LT,
         [("La continuité", " : un an sur le même produit, là où les projets sont morcelés."),
          ("La responsabilité", " : le code engage le projet, l'équipe, l'entreprise."),
          ("La vision produit", " : déploiement, pannes, tests — ancrés par le réel.")]),
        ("↑  Ce que j'ai dû renforcer", PINK_DK, PINK_LT,
         [("L'IA comme outil", " de dev : peu abordée en formation (logique d'évaluation)."),
          ("L'arbitrage sécurité / fonctionnel", " : entrevu chez M. Chikhaoui, vécu en entreprise."),
          ("Le relationnel à l'international", " : échanges en anglais (Inde) — un axe de progrès assumé.")])]
x0=Inches(0.7); cw=Inches(3.83); gap=Inches(0.2); y=Inches(2.15)
for i,(h,hc,fl,items) in enumerate(cols14):
    column_card(s, x0+i*(cw+gap), y, cw, Inches(3.35), h, hc, items, marker="—", fill=fl)
takeaway(s, Inches(0.7), Inches(5.75), Inches(11.9), "L'IA, des deux côtés  ",
         "outil de développement ET matière du produit (elle y propose savoirs et HSC) — avoir d'abord appris à coder sans elle m'a permis de la mobiliser, pas de m'y substituer.", BASE)

# =========================================================================
# SLIDE 15 — BÉNÉFICES & ÉVOLUTION
# =========================================================================
s = new_slide(BASE, 15)
title_block(s, "PARTIE III · BÉNÉFICES", "Pour l'entreprise, pour le projet — et pour moi", BASE)
column_card(s, Inches(0.7), Inches(2.15), Inches(5.85), Inches(2.75), "Pour l'A.F.D.E.C & le projet", INDIGO,
            [("Stratégique", " : Docker → Cloud Run, on paie l'usage. Une frugalité qui rend le projet soutenable."),
             ("Patrimonial", " : un actif documenté, testé, déployé — non plus une démo dépendante de son auteur."),
             ("Pensé commercial", " : la base embarque ses données de test, le multi-entités dès la conception.")],
            marker="—", fill=PANEL2)
column_card(s, Inches(6.75), Inches(2.15), Inches(5.85), Inches(2.75), "Pour moi", PINK_DK,
            [("Montée en responsabilité", " : de stagiaire à une forme d'association au projet."),
             ("De l'exécution à la co-construction", " : défendre une architecture, accepter qu'elle évolue."),
             ("Maturité d'ingénieur produit", " : du « faire marcher » au « faire tenir ».")],
            marker="—", fill=PINK_LT)
rounded(s, Inches(0.7), Inches(5.2), Inches(11.9), Inches(1.1), fill=INDIGO_LT, lc=None, rad=0.07)
rect(s, Inches(0.7), Inches(5.2), Inches(0.1), Inches(1.1), fill=INDIGO)
text(s, Inches(1.1), Inches(5.2), Inches(11.2), Inches(1.1),
     [[("« En apprenant à faire grandir un projet, j'ai grandi avec lui. »", {'font':TITLE,'color':INDIGO,'size':22,'bold':True})]],
     anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)

# =========================================================================
# SLIDE 16 — CONCLUSION
# =========================================================================
s = new_slide(BASE, 16)
title_block(s, "CONCLUSION", "J'avais repris une maquette ; je laisse un produit", BASE)
bullets(s, Inches(0.72), Inches(2.25), Inches(11.8), Inches(2.2),
        [("La technique n'est jamais neutre", " : derrière un filtre entity_id, le respect des données d'autrui ; derrière un fallback, le refus de laisser tomber l'utilisateur."),
         ("Nommer ce qui n'est pas fini", " — sécurité, optimisation, distribution — non comme des échecs, mais comme le trajet qui reste."),
         ("Le centre de gravité se déplace", " du développement vers la commercialisation : un terrain nouveau, et passionnant.")],
        marker="—", mcolor=INDIGO, size=14.5, sp_after=13, ls=1.12)
rounded(s, Inches(0.7), Inches(5.0), Inches(11.9), Inches(1.35), fill=GREEN_LT, lc=None, rad=0.06)
rect(s, Inches(0.7), Inches(5.0), Inches(0.1), Inches(1.35), fill=GREEN_DK)
text(s, Inches(1.1), Inches(5.12), Inches(11.2), Inches(0.5),
     [[("La réponse à la problématique  ", {'color':GREEN_DK,'size':13.5,'bold':True}),
       ("transformer une maquette en produit fut bien un véritable levier de montée en compétences.", {'color':INK,'size':13.5})]], ls=1.05)
text(s, Inches(1.1), Inches(5.66), Inches(11.2), Inches(0.6),
     [[("La méthode OPTIQ veut révéler la place de chacun dans une organisation ; en la rendant accessible, j'ai surtout découvert la mienne.",
        {'color':SLATE,'size':13,'italic':True})]], ls=1.08)

# =========================================================================
# SLIDE 17 — MERCI / QUESTIONS
# =========================================================================
s = new_slide(BASE, 17)
# brand bands
rect(s, 0,0, SW, Inches(0.16), fill=PINK)
rect(s, 0, Inches(0.16), SW, Inches(0.05), fill=GREEN)
text(s, Inches(0.6), Inches(2.5), Inches(12.1), Inches(1.2),
     [[("Merci de votre attention.", {'font':TITLE,'color':INK,'size':46,'bold':True})]], align=PP_ALIGN.CENTER)
rect(s, Inches(5.96), Inches(3.95), Inches(1.4), Pt(3.2), fill=PINK)
rect(s, Inches(7.42), Inches(3.9), Inches(0.15), Inches(0.15), fill=GREEN)
text(s, Inches(0.6), Inches(4.25), Inches(12.1), Inches(0.6),
     [[("Place à vos questions.", {'color':SLATE,'size':19,'italic':True})]], align=PP_ALIGN.CENTER)
img_fit(s, asset("img01.png"), Inches(4.55), Inches(5.5), Inches(2.0), Inches(0.85), frame=False)
img_fit(s, asset("img00.png"), Inches(6.95), Inches(5.5), Inches(1.85), Inches(0.8), frame=False)

prs.save("Soutenance_DevOPTIQ_Mael_Girardin.pptx")
print("Saved:", os.path.abspath("Soutenance_DevOPTIQ_Mael_Girardin.pptx"))
print("Slides:", len(prs.slides._sldIdLst))
