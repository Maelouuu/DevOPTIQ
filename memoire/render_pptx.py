# -*- coding: utf-8 -*-
"""Rendu PIL fidèle (géométrie/couleurs/images) du pptx généré, pour QA visuelle."""
import re, io, sys
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

EMU=914400; SCALE=132.0
def px(emu): return int(round((emu or 0)/EMU*SCALE))
def ptpx(pt): return pt*SCALE/72.0

FONTS={}
FB="/usr/share/fonts/truetype/liberation/"
def fontfile(name, bold, italic):
    n=(name or "").lower()
    if "consol" in n or "mono" in n: base="LiberationMono"
    elif "georgia" in n or "serif" in n or "times" in n or "fraunces" in n: base="LiberationSerif"
    else: base="LiberationSans"
    if bold and italic: suf="-BoldItalic"
    elif bold: suf="-Bold"
    elif italic: suf="-Italic"
    else: suf="-Regular"
    return FB+base+suf+".ttf"
def get_font(name,size_px,bold,italic):
    key=(name,size_px,bool(bold),bool(italic))
    if key not in FONTS:
        FONTS[key]=ImageFont.truetype(fontfile(name,bold,italic), max(6,int(round(size_px))))
    return FONTS[key]

def rgb_of(c):
    h=str(c); return (int(h[0:2],16),int(h[2:4],16),int(h[4:6],16))
def shape_fill(sp):
    try:
        from pptx.enum.dml import MSO_FILL
        if sp.fill.type==MSO_FILL.SOLID: return rgb_of(sp.fill.fore_color.rgb)
    except Exception: pass
    return None
def shape_line(sp):
    try:
        c=sp.line.color.rgb
        if c is None: return None,0
        w=sp.line.width; wpx=px(w) if w else max(1,int(ptpx(1)))
        return rgb_of(c), max(1,wpx)
    except Exception:
        return None,0

def draw_shape_bg(d, sp):
    x,y,w,h=px(sp.left),px(sp.top),px(sp.width),px(sp.height)
    fill=shape_fill(sp); lc,lw=shape_line(sp)
    if fill is None and lc is None: return
    typ=None
    try: typ=sp.auto_shape_type
    except Exception: typ=None
    name=str(typ)
    if "OVAL" in name:
        d.ellipse([x,y,x+w,y+h], fill=fill, outline=lc, width=lw or 1)
    elif "ROUNDED" in name:
        rad=int(min(w,h)*0.12)
        try:
            adj=sp.adjustments[0]; rad=int(min(w,h)*float(adj))
        except Exception: pass
        rad=max(2,min(rad,min(w,h)//2))
        d.rounded_rectangle([x,y,x+w,y+h], radius=rad, fill=fill, outline=lc, width=lw or 1)
    else:
        d.rectangle([x,y,x+w,y+h], fill=fill, outline=lc, width=lw or 1)

def run_color(run):
    try:
        c=run.font.color.rgb
        if c is not None: return rgb_of(c)
    except Exception: pass
    return (15,23,42)

def layout_para(d, runs, max_w):
    words=[]
    for txt,font,color in runs:
        for part in re.split(r'(\s+)', txt):
            if part=='' : continue
            words.append((part,font,color))
    lines=[]; cur=[]; cur_w=0
    for tok,font,color in words:
        w=d.textlength(tok,font=font)
        if tok.isspace():
            if cur: cur.append([tok,font,color,w]); cur_w+=w
            continue
        if cur_w+w>max_w and cur:
            while cur and cur[-1][0].isspace(): cur_w-=cur[-1][3]; cur.pop()
            lines.append([cur,cur_w]); cur=[]; cur_w=0
        cur.append([tok,font,color,w]); cur_w+=w
    if cur:
        while cur and cur[-1][0].isspace(): cur_w-=cur[-1][3]; cur.pop()
        lines.append([cur,cur_w])
    return lines

def draw_text(d, sp):
    tf=sp.text_frame
    x,y,w,h=px(sp.left),px(sp.top),px(sp.width),px(sp.height)
    ml=px(tf.margin_left); mr=px(tf.margin_right); mt=px(tf.margin_top); mb=px(tf.margin_bottom)
    bx=x+ml; by=y+mt; bw=w-ml-mr; bh=h-mt-mb
    # build paragraphs
    paras=[]
    for p in tf.paragraphs:
        runs=[]
        for r in p.runs:
            sz=r.font.size; pt=sz.pt if sz is not None else 14
            f=get_font(r.font.name, ptpx(pt), r.font.bold, r.font.italic)
            runs.append((r.text, f, run_color(r)))
        if not runs: runs=[("",get_font("",ptpx(14),False,False),(15,23,42))]
        ls=p.line_spacing or 1.04
        sa=p.space_after.pt if p.space_after is not None else 4
        al=p.alignment
        lines=layout_para(d, runs, bw)
        # line height from max font in each line
        lh=[]
        for toks,_w in lines:
            mx=max((t[1].size for t in toks), default=int(ptpx(14)))
            lh.append(mx*ls*1.16)
        paras.append((lines,lh,sa,al))
    total=sum(sum(lh)+ (sa*SCALE/72.0) for lines,lh,sa,al in paras)
    # vertical anchor
    va=tf.vertical_anchor
    if va==MSO_ANCHOR.MIDDLE: cy=by+max(0,(bh-total)/2)
    elif va==MSO_ANCHOR.BOTTOM: cy=by+max(0,bh-total)
    else: cy=by
    for lines,lh,sa,al in paras:
        for i,(toks,lw_) in enumerate(lines):
            if al==PP_ALIGN.CENTER: lx=bx+(bw-lw_)/2
            elif al==PP_ALIGN.RIGHT: lx=bx+(bw-lw_)
            else: lx=bx
            # baseline-ish: draw at cy, vertical center within lh[i]
            asc=max((t[1].size for t in toks), default=int(ptpx(14)))
            ty=cy+(lh[i]-asc)/2 - asc*0.06
            for tok,font,color,wtok in toks:
                d.text((lx,ty),tok,font=font,fill=color)
                lx+=wtok
            cy+=lh[i]
        cy+=sa*SCALE/72.0

def render_slide(slide, idx):
    W=px(prs.slide_width); H=px(prs.slide_height)
    # background
    bg=(255,255,255)
    try:
        from pptx.enum.dml import MSO_FILL
        if slide.background.fill.type==MSO_FILL.SOLID:
            bg=rgb_of(slide.background.fill.fore_color.rgb)
    except Exception: pass
    img=Image.new("RGB",(W,H),bg); d=ImageDraw.Draw(img)
    for sp in slide.shapes:
        if sp.shape_type==MSO_SHAPE_TYPE.PICTURE:
            try:
                blob=sp.image.blob; pim=Image.open(io.BytesIO(blob)).convert("RGB")
                pim=pim.resize((max(1,px(sp.width)),max(1,px(sp.height))))
                img.paste(pim,(px(sp.left),px(sp.top)))
            except Exception as e: print("img err",e)
            continue
        if sp.has_text_frame:
            draw_shape_bg(d,sp); draw_text(d,sp)
        else:
            draw_shape_bg(d,sp)
    return img

prs=Presentation("Soutenance_DevOPTIQ_Mael_Girardin.pptx")
imgs=[]
for i,slide in enumerate(prs.slides):
    im=render_slide(slide,i); im.save(f"_slide{i+1:02d}.png"); imgs.append(im)
print("rendered",len(imgs),"slides")
# contact sheet
cols=3; tw=imgs[0].width//3; th=imgs[0].height//3
sheet=Image.new("RGB",(cols*tw, ((len(imgs)+cols-1)//cols)*th),(230,230,230))
dd=ImageDraw.Draw(sheet)
for i,im in enumerate(imgs):
    t=im.resize((tw-6,th-6)); x=(i%cols)*tw; y=(i//cols)*th
    sheet.paste(t,(x+3,y+3)); dd.text((x+6,y+4),str(i+1),fill=(190,24,93))
sheet.save("_deck_contact.png"); print("contact ->",sheet.size)
